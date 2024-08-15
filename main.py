from pptx import Presentation
from dotenv import load_dotenv
import os
from io import BytesIO

load_dotenv()

pptx_path = os.getenv("PPTX_PATH")

def copy_slide(new_pptx, ex_slide):
    """
        Función que copia un slide y lo pega en una presentación de Power Point.
        Args:
            new_pptx (pptx.presentation.Presentation): La presentación de Power Point a la que se desea agregar el slide.
            ex_slide (pptx.slide.Slide): El slide que se desea copiar y pegar en la presentación.
        Returns:
            pptx.slide.Slide: El slide copiado.
    """
    new_slide = new_pptx.slides.add_slide(ex_slide.slide_layout) # Se crea un nuevo slide desde cero.
    # Se evalúan todos los shapes del slide a copiar.
    for shape in ex_slide.shapes:
        # Si el shape actual es un placeholder.
        if shape.is_placeholder:
            # Se crea un nuevo placeholder.
            phf = new_slide.placeholders[shape.placeholder_format.idx]
            # Se agrega el texto al nuevo placeholder en caso de existir.
            if shape.has_text_frame:
                phf.text = shape.text_frame.text
            # Se agrega el chart al nuevo placeholder en caso de existir.
            elif shape.has_chart:
                chart = shape.chart
                phf.insert_chart(
                    chart.chart_type,
                    chart_data=chart.chart_data
                )
            # Se agrega la tabla al nuevo placeholder en caso de existir.
            elif shape.has_table:
                table = shape.table
                # Se debe copiar el texto de cada una de las celdas de la tabla.
                for i, row in enumerate(table.rows):
                    for j, cell in enumerate(row.cells):
                        phf.table.cell(i, j).text = cell.text
            # Se agrega la imagen al nuevo placeholder en caso de existir.
            elif shape.has_picture:
                image = shape.image
                phf.insert_picture(image)
        elif shape.shape_type == 17:  # Si el shape actual es un cuadro de texto.
            new_shape = new_slide.shapes.add_textbox(
                shape.left, shape.top, shape.width, shape.height
            )
            new_shape.text_frame.text = shape.text_frame.text
        elif shape.shape_type == 13:  # Si el shape actual es una imagen.
            image = shape.image
            new_slide.shapes.add_picture(BytesIO(image.blob), shape.left, shape.top, shape.width, shape.height)
        else: # Para los otros casos.
            # Si el shape actual es un chart.
            if shape.has_chart:
                chart = shape.chart
                new_chart = new_slide.shapes.add_chart(
                    chart.chart_type, 
                    shape.left, shape.top, shape.width, shape.height, 
                    chart_data=chart.chart_data
                ).chart
                
                new_chart.has_legend = chart.has_legend
                new_chart.legend.position = chart.legend.position

            elif shape.has_table: # Si el shape actual es una tabla.
                table = shape.table
                new_table = new_slide.shapes.add_table(
                    rows=len(table.rows),
                    cols=len(table.columns),
                    left=shape.left,
                    top=shape.top,
                    width=shape.width,
                    height=shape.height
                ).table
                # Se copia el contenido de cada una de las celdas.
                for i, row in enumerate(table.rows):
                    for j, cell in enumerate(row.cells):
                        new_table.cell(i, j).text = cell.text

            else: # Si el shape actual es un shape automático.
                new_shape = new_slide.shapes.add_shape(
                    shape.shape_type,
                    shape.left, shape.top, shape.width, shape.height
                )
            # Agregar texto a las opciones anteriores en caso de existir.
            if shape.has_text_frame:
                new_shape.text = shape.text_frame.text

    return new_slide

def copy_pptx():
    new_pptx = Presentation() # La nueva presentación
    for root, _, files in os.walk(pptx_path):
        for file in files:
            act_pptx = Presentation(f'{root}/{file}')
            for slide in act_pptx.slides:
                copy_slide(new_pptx, slide)

    try:
        os.makedirs('output/')
    except FileExistsError:
        pass
    
    new_pptx.save('output/output.pptx')
    os.system("cls")
    print("SE HAN JUNTADO LAS PRESENTACIONES CORRECTAMENTE!")
    os.system("pause")
copy_pptx()